#!/usr/bin/env python3
"""Compare report files by rendered pixels and/or PowerPoint objects.

Python version: 3.9+

Supported pixel inputs:
- Images supported by Pillow: PNG, JPG, JPEG, BMP, TIFF, WEBP.
- PPT/PPTX files on Windows when Microsoft PowerPoint + pywin32 are installed.

Supported object inputs:
- PPTX files when python-pptx is installed.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import sys
import tempfile
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

from PIL import Image, ImageChops, ImageDraw, ImageFilter


PPT_EXTS = {".ppt", ".pptx"}
PPTX_EXT = ".pptx"
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
EMU_PER_POINT = 12700
DEFAULT_OUTPUT_DIR = "output"
DEFAULT_GAP_PENALTY = -0.12
SIGNATURE_SIZE = (128, 72)
BLOCK_GRID = (8, 6)
COLOR_BINS = 8
WHITE_THRESHOLD = 240  # gray >= this counts as blank background, below = slide content
EDGE_THRESHOLD = 32  # edge response above this counts as a structural edge

# Tokens that vary between report versions (data values, identifiers, dates) but
# don't change the slide's template / role. Stripping them before text matching
# keeps the comparison focused on the boilerplate that *defines* the template.
_VARIABLE_TOKEN_PATTERN = re.compile(
    r"\d{1,4}[/\-.]\d{1,2}[/\-.]\d{1,4}"          # dates: 2026/03/25, 2025-9-17
    r"|[A-Za-z]?\d+(?:[A-Za-z]+\d*)*"             # codes: T160, K14C, PT2, 4000rpm
    r"|\d+(?:\.\d+)?"                              # numbers: 4000, 106.1, 24.0
)


@dataclass
class PageDiff:
    page: int
    expected_page: Optional[int]
    actual_page: Optional[int]
    match_score: Optional[float]
    match_status: str
    width: int
    height: int
    compared_pixels: int
    different_pixels: int
    difference_percent: float
    max_channel_delta: int
    bbox: Optional[Tuple[int, int, int, int]]
    passed: bool
    output_overlay: str
    output_mask: str
    output_expected: str
    output_actual: str
    regions: List[Tuple[int, int, int, int]]


@dataclass
class SlideMatch:
    expected_index: Optional[int]
    actual_index: Optional[int]
    score: Optional[float]
    status: str


@dataclass
class ObjectDiff:
    pair_index: int                  # 1-based pair number, matches PageDiff.page
    expected_slide: Optional[int]    # 1-based slide number in the expected report
    actual_slide: Optional[int]      # 1-based slide number in the actual report
    object_index: int
    field: str
    expected: Any
    actual: Any


@dataclass
class SlideSignature:
    """Cached per-slide features, reused across the N×M alignment comparison.

    Every feature here is *content-focused* — it ignores the blank background so
    two unrelated slides that happen to share lots of whitespace are not scored
    as similar.
    """

    content_mask: Image.Image       # binary: where the slide has non-white content
    edge_mask: Image.Image          # binary: structural edges (chart frames, dividers)
    block_ink: List[float]          # ink ratio in [0, 1] per grid cell
    content_color_hist: List[float]  # colour distribution over content pixels only
    text: Optional[str]             # full slide text (PPTX only)
    title: Optional[str] = None     # title placeholder text (PPTX only)


def parse_color(value: str) -> Tuple[int, int, int]:
    parts = value.split(",")
    if len(parts) != 3:
        raise argparse.ArgumentTypeError("Color must be formatted as R,G,B")

    try:
        rgb = tuple(int(part.strip()) for part in parts)
    except ValueError as exc:
        raise argparse.ArgumentTypeError("Color values must be integers") from exc

    if any(channel < 0 or channel > 255 for channel in rgb):
        raise argparse.ArgumentTypeError("Color values must be between 0 and 255")

    return rgb  # type: ignore[return-value]


def safe_text(value: Any) -> str:
    return "" if value is None else str(value)


def export_ppt_to_png(path: Path, dpi: int, temp_root: Path) -> List[Path]:
    try:
        import win32com.client  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "PPT/PPTX pixel comparison requires pywin32. Install it with: pip install pywin32"
        ) from exc

    export_dir = temp_root / path.stem
    export_dir.mkdir(parents=True, exist_ok=True)

    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = None
    try:
        presentation = powerpoint.Presentations.Open(
            str(path.resolve()),
            WithWindow=False,
            ReadOnly=True,
        )
        width = int(round((presentation.PageSetup.SlideWidth / 72.0) * dpi))
        height = int(round((presentation.PageSetup.SlideHeight / 72.0) * dpi))
        presentation.Export(str(export_dir), "PNG", width, height)
    finally:
        if presentation is not None:
            presentation.Close()
        powerpoint.Quit()

    exported = sorted(export_dir.glob("*.PNG"), key=slide_export_sort_key)
    if not exported:
        exported = sorted(export_dir.glob("*.png"), key=slide_export_sort_key)
    if not exported:
        raise RuntimeError("PowerPoint did not export any slide images.")
    return exported


def slide_export_sort_key(path: Path) -> Tuple[int, str]:
    digits = "".join(character for character in path.stem if character.isdigit())
    number = int(digits) if digits else 0
    return number, path.name.lower()


def load_ppt_pages(path: Path, dpi: int, temp_root: Path) -> List[Image.Image]:
    image_paths = export_ppt_to_png(path, dpi, temp_root)
    pages = []
    for image_path in image_paths:
        with Image.open(image_path) as image:
            pages.append(image.convert("RGB"))
    return pages


def load_pages(path: Path, dpi: int, temp_root: Path) -> List[Image.Image]:
    suffix = path.suffix.lower()
    if suffix in PPT_EXTS:
        return load_ppt_pages(path, dpi, temp_root)
    if suffix not in IMAGE_EXTS:
        raise ValueError(
            "Unsupported input type '{}'. Use an image, PPT, or PPTX file.".format(
                path.suffix
            )
        )

    with Image.open(path) as image:
        return [image.convert("RGB")]


def fit_to_same_canvas(
    expected: Image.Image,
    actual: Image.Image,
    background: Tuple[int, int, int] = (255, 255, 255),
) -> Tuple[Image.Image, Image.Image]:
    width = max(expected.width, actual.width)
    height = max(expected.height, actual.height)

    expected_canvas = Image.new("RGB", (width, height), background)
    actual_canvas = Image.new("RGB", (width, height), background)
    expected_canvas.paste(expected.convert("RGB"), (0, 0))
    actual_canvas.paste(actual.convert("RGB"), (0, 0))

    return expected_canvas, actual_canvas


def build_threshold_mask(diff: Image.Image, threshold: int) -> Image.Image:
    channel_masks = [
        channel.point(lambda value: 255 if value > threshold else 0)
        for channel in diff.convert("RGB").split()
    ]
    return ImageChops.lighter(ImageChops.lighter(channel_masks[0], channel_masks[1]), channel_masks[2])


def count_mask_pixels(mask: Image.Image) -> int:
    histogram = mask.histogram()
    return sum(count for value, count in enumerate(histogram) if value > 0)


def mask_to_regions(mask: Image.Image, min_area: int = 16, padding: int = 3) -> List[Tuple[int, int, int, int]]:
    binary = mask.convert("1")
    width, height = binary.size
    pixels = binary.load()
    visited = set()
    regions = []

    for start_y in range(height):
        for start_x in range(width):
            point = (start_x, start_y)
            if point in visited or not pixels[start_x, start_y]:
                continue

            stack = [point]
            visited.add(point)
            min_x = max_x = start_x
            min_y = max_y = start_y
            area = 0

            while stack:
                x, y = stack.pop()
                area += 1
                min_x = min(min_x, x)
                max_x = max(max_x, x)
                min_y = min(min_y, y)
                max_y = max(max_y, y)

                for nx, ny in ((x - 1, y), (x + 1, y), (x, y - 1), (x, y + 1)):
                    if nx < 0 or ny < 0 or nx >= width or ny >= height:
                        continue
                    neighbor = (nx, ny)
                    if neighbor in visited or not pixels[nx, ny]:
                        continue
                    visited.add(neighbor)
                    stack.append(neighbor)

            if area >= min_area:
                regions.append(
                    (
                        max(0, min_x - padding),
                        max(0, min_y - padding),
                        min(width, max_x + 1 + padding),
                        min(height, max_y + 1 + padding),
                    )
                )

    return regions


def make_overlay(
    base: Image.Image,
    regions: Sequence[Tuple[int, int, int, int]],
    highlight_color: Tuple[int, int, int],
    alpha: int,
) -> Image.Image:
    overlay = base.convert("RGBA")
    highlight_layer = Image.new("RGBA", overlay.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(highlight_layer)
    outline = (highlight_color[0], highlight_color[1], highlight_color[2], 255)
    image_area = max(1, overlay.size[0] * overlay.size[1])

    for region in regions:
        x1, y1, x2, y2 = region
        region_area = max(1, (x2 - x1) * (y2 - y1))
        fill_alpha = alpha
        if region_area / float(image_area) >= 0.08:
            fill_alpha = max(8, alpha // 5)
        fill = (highlight_color[0], highlight_color[1], highlight_color[2], fill_alpha)
        draw.rectangle(region, fill=fill, outline=outline)
        for inset in range(1, 4):
            if x2 - x1 > inset * 2 and y2 - y1 > inset * 2:
                draw.rectangle((x1 + inset, y1 + inset, x2 - inset, y2 - inset), outline=outline)

    return Image.alpha_composite(overlay, highlight_layer).convert("RGB")


def compare_page(
    expected: Image.Image,
    actual: Image.Image,
    page_number: int,
    expected_page: Optional[int],
    actual_page: Optional[int],
    match_score: Optional[float],
    match_status: str,
    output_dir: Path,
    threshold: int,
    allowed_percent: float,
    highlight_color: Tuple[int, int, int],
    alpha: int,
) -> PageDiff:
    expected_canvas, actual_canvas = fit_to_same_canvas(expected, actual)
    diff = ImageChops.difference(expected_canvas, actual_canvas)
    mask = build_threshold_mask(diff, threshold)

    different_pixels = count_mask_pixels(mask)
    compared_pixels = mask.width * mask.height
    difference_percent = (different_pixels / compared_pixels) * 100 if compared_pixels else 0
    extrema = diff.getextrema()
    max_channel_delta = max(channel_max for _channel_min, channel_max in extrema)
    bbox = mask.getbbox()
    passed = difference_percent <= allowed_percent
    regions = mask_to_regions(mask)

    overlay = make_overlay(actual_canvas, regions, highlight_color, alpha)
    expected_label = expected_page if expected_page is not None else 0
    actual_label = actual_page if actual_page is not None else 0
    file_stem = "match_{:03d}_expected_{:03d}_actual_{:03d}".format(page_number, expected_label, actual_label)
    overlay_path = output_dir / "{}_overlay.png".format(file_stem)
    mask_path = output_dir / "{}_mask.png".format(file_stem)
    expected_path = output_dir / "{}_expected.png".format(file_stem)
    actual_path = output_dir / "{}_actual.png".format(file_stem)
    overlay.save(overlay_path)
    mask.save(mask_path)
    expected_canvas.save(expected_path)
    actual_canvas.save(actual_path)

    return PageDiff(
        page=page_number,
        expected_page=expected_page,
        actual_page=actual_page,
        match_score=round(match_score, 6) if match_score is not None else None,
        match_status=match_status,
        width=mask.width,
        height=mask.height,
        compared_pixels=compared_pixels,
        different_pixels=different_pixels,
        difference_percent=round(difference_percent, 6),
        max_channel_delta=max_channel_delta,
        bbox=bbox,
        passed=passed,
        output_overlay=str(overlay_path),
        output_mask=str(mask_path),
        output_expected=str(expected_path),
        output_actual=str(actual_path),
        regions=regions,
    )



def _content_mask(image: Image.Image, size: Tuple[int, int] = SIGNATURE_SIZE) -> Image.Image:
    """Binary mask (0/255) marking where the slide carries content vs. blank paper.

    Report slides are mostly white. If we compare whole thumbnails, the matching
    white background drowns out the few percent of pixels that actually differ,
    so unrelated slides score deceptively high. Reducing each slide to *where the
    ink is* makes the comparison about content, not paper.  A light dilation adds
    tolerance to a few pixels of render jitter.
    """
    gray = image.convert("L").resize(size, Image.BILINEAR)
    mask = gray.point(lambda value: 255 if value < WHITE_THRESHOLD else 0)
    return mask.filter(ImageFilter.MaxFilter(3))


def _edge_mask(image: Image.Image, size: Tuple[int, int] = SIGNATURE_SIZE) -> Image.Image:
    """Binary mask of structural edges — chart frames, table borders, dividers.

    FIND_EDGES runs on the full-resolution image (before the downscale) so thin
    boundaries survive; the result is thresholded to a clean on/off edge map.
    """
    edges = image.convert("L").filter(ImageFilter.FIND_EDGES).resize(size, Image.BILINEAR)
    mask = edges.point(lambda value: 255 if value > EDGE_THRESHOLD else 0)
    return mask.filter(ImageFilter.MaxFilter(3))


def _mask_iou(mask_a: Image.Image, mask_b: Image.Image) -> float:
    """Intersection-over-union of two binary masks — measures spatial overlap.

    This is the core layout signal: content in the same places => high IoU;
    content arranged differently => low IoU, regardless of how much whitespace
    the two slides share.
    """
    intersection = ImageChops.darker(mask_a, mask_b)
    union = ImageChops.lighter(mask_a, mask_b)
    intersection_count = sum(c for v, c in enumerate(intersection.histogram()) if v > 0)
    union_count = sum(c for v, c in enumerate(union.histogram()) if v > 0)
    if union_count == 0:
        return 1.0  # both slides are blank
    return intersection_count / union_count


def _block_ink_ratios(image: Image.Image, grid: Tuple[int, int] = BLOCK_GRID) -> List[float]:
    """Fraction of content (non-white) pixels in each grid cell, range [0, 1].

    Unlike mean brightness — where a 95%-white cell and a 90%-white cell differ
    by a negligible 4% — ink ratio spans the full [0, 1] range, so a cell with a
    dense chart reads very differently from a near-empty one.
    """
    gray = image.convert("L")
    cols, rows = grid
    width, height = gray.size
    ratios: List[float] = []
    for row in range(rows):
        for col in range(cols):
            box = (
                width * col // cols,
                height * row // rows,
                width * (col + 1) // cols,
                height * (row + 1) // rows,
            )
            histogram = gray.crop(box).histogram()
            total = sum(histogram) or 1
            ink = sum(c for v, c in enumerate(histogram) if v < WHITE_THRESHOLD)
            ratios.append(ink / total)
    return ratios


def _block_ink_similarity(ratios_a: Sequence[float], ratios_b: Sequence[float]) -> float:
    if not ratios_a or not ratios_b or len(ratios_a) != len(ratios_b):
        return 0.0
    mean_diff = sum(abs(a - b) for a, b in zip(ratios_a, ratios_b)) / len(ratios_a)
    return max(0.0, 1.0 - mean_diff)


def _content_color_histogram(image: Image.Image, bins: int = COLOR_BINS) -> List[float]:
    """Quantized RGB histogram over *content* pixels only — the white background
    is excluded so the distribution reflects what is actually drawn.

    A slide of red/green/yellow polar curves then has a genuinely different
    colour profile from a blue line graph, instead of both being "mostly white".
    """
    rgb = image.convert("RGB").resize((96, 54), Image.BILINEAR)
    bucket = 256 // bins
    hist = [0] * (bins * 3)
    for r, g, b in rgb.getdata():
        if r >= WHITE_THRESHOLD and g >= WHITE_THRESHOLD and b >= WHITE_THRESHOLD:
            continue
        hist[min(bins - 1, r // bucket)] += 1
        hist[bins + min(bins - 1, g // bucket)] += 1
        hist[bins * 2 + min(bins - 1, b // bucket)] += 1
    total = sum(hist)
    if total == 0:
        return [0.0] * (bins * 3)
    return [value / total for value in hist]


def _histogram_intersection(hist_a: Sequence[float], hist_b: Sequence[float]) -> float:
    if len(hist_a) != len(hist_b):
        return 0.0
    return sum(min(a, b) for a, b in zip(hist_a, hist_b))


def _strip_variable_tokens(text: str) -> str:
    """Remove data values, identifiers, and dates — the parts that change between
    report versions while the slide template stays the same.

    Two slides with the same title "解析条件 - 燃焼圧" but different rpm values
    (4000/5500/6000 vs 6000) should read as nearly-identical text after this
    pass; their template tokens (解析条件, 燃焼圧, 筒内圧, クランク角, 仕様, …)
    remain intact.
    """
    return _VARIABLE_TOKEN_PATTERN.sub(" ", text)


def _char_bigrams(text: str) -> set:
    """Character bigrams over the slide's template tokens — works for CJK text
    where whitespace tokenizing fails, and ignores variable data values."""
    compact = "".join(_strip_variable_tokens(text).split()).lower()
    if len(compact) < 2:
        return {compact} if compact else set()
    return {compact[index:index + 2] for index in range(len(compact) - 1)}


def text_similarity(text_a: str, text_b: str) -> float:
    """Jaccard similarity over character bigrams of two slides' text content."""
    bigrams_a = _char_bigrams(text_a or "")
    bigrams_b = _char_bigrams(text_b or "")
    if not bigrams_a and not bigrams_b:
        return 1.0
    if not bigrams_a or not bigrams_b:
        return 0.0
    union = len(bigrams_a | bigrams_b)
    return len(bigrams_a & bigrams_b) / union if union else 0.0


def compute_signature(
    image: Image.Image,
    text: Optional[str] = None,
    title: Optional[str] = None,
) -> SlideSignature:
    """Pre-compute every per-slide feature once so alignment can reuse them.

    align_slide_pages compares every expected slide against every actual slide
    (N×M); without caching, each slide's signatures would be recomputed N or M
    times over.
    """
    return SlideSignature(
        content_mask=_content_mask(image),
        edge_mask=_edge_mask(image),
        block_ink=_block_ink_ratios(image),
        content_color_hist=_content_color_histogram(image),
        text=text,
        title=title,
    )


def signature_similarity(expected: SlideSignature, actual: SlideSignature) -> float:
    """Composite similarity in [0, 1] from all available per-slide signals.

    Every visual component looks at *content*, never at the shared white
    background, so the score reflects how alike the slides actually are:
      40% layout — IoU of content masks (is the content in the same places?)
      25% blocks — per-cell ink ratio (how is content distributed?)
      20% edges  — IoU of structural edges (chart frames, dividers, tables)
      15% colour — colour distribution over content pixels

    The mix with text depends on whether titles are available:

    - When **both** slides expose a title (PPTX title placeholder), the title is
      the strongest template signal in technical reports — two slides titled
      "解析条件 - 燃焼圧" almost certainly play the same role even if their data
      tables differ in row count. If those titles match strongly the visual is
      mostly a tie-breaker; otherwise the four signals (visual + title + body +
      a soft body match) blend.
    - When only full text is available, fall back to a balanced visual/text mix.
    - With no text at all (images, PPT) the visual score stands on its own.
    """
    layout_sim = _mask_iou(expected.content_mask, actual.content_mask)
    block_sim = _block_ink_similarity(expected.block_ink, actual.block_ink)
    edge_sim = _mask_iou(expected.edge_mask, actual.edge_mask)
    color_sim = _histogram_intersection(expected.content_color_hist, actual.content_color_hist)
    visual = 0.40 * layout_sim + 0.25 * block_sim + 0.20 * edge_sim + 0.15 * color_sim

    has_title = bool(expected.title) and bool(actual.title)
    has_text = expected.text is not None and actual.text is not None

    if has_title:
        title_sim = text_similarity(expected.title or "", actual.title or "")
        body_sim = text_similarity(expected.text or "", actual.text or "") if has_text else title_sim
        text_score = 0.70 * title_sim + 0.30 * body_sim
        if title_sim >= 0.85:
            # Titles are essentially identical → strong template match. Trust the
            # text signal; visual is just a tie-breaker against false positives.
            return 0.15 * visual + 0.85 * text_score
        if title_sim >= 0.55:
            # Related titles — meaningful template overlap.
            return 0.35 * visual + 0.65 * text_score
        # Title mismatch but body might still help; treat like plain text mode.
        return 0.55 * visual + 0.45 * text_score

    if has_text:
        return 0.55 * visual + 0.45 * text_similarity(expected.text, actual.text)
    return visual


def slide_similarity(expected: Image.Image, actual: Image.Image) -> float:
    """Convenience wrapper: composite visual similarity for two raw images."""
    return signature_similarity(compute_signature(expected), compute_signature(actual))


def build_score_matrix(
    expected_sigs: Sequence[SlideSignature],
    actual_sigs: Sequence[SlideSignature],
) -> List[List[float]]:
    """Full N×M similarity matrix between every expected and actual slide."""
    return [
        [signature_similarity(expected_sig, actual_sig) for actual_sig in actual_sigs]
        for expected_sig in expected_sigs
    ]


def align_from_score_matrix(
    scores: Sequence[Sequence[float]],
    expected_count: int,
    actual_count: int,
    min_match_score: float,
    gap_penalty: float = DEFAULT_GAP_PENALTY,
) -> List[SlideMatch]:
    """Needleman–Wunsch style alignment over a precomputed similarity matrix.

    Building the matrix once and feeding it here keeps scoring and alignment
    separate, so the same matrix can also be exported for the user to audit.
    """
    if expected_count == 0 and actual_count == 0:
        return []

    dp = [[0.0] * (actual_count + 1) for _ in range(expected_count + 1)]
    step = [[""] * (actual_count + 1) for _ in range(expected_count + 1)]

    for i in range(1, expected_count + 1):
        dp[i][0] = dp[i - 1][0] + gap_penalty
        step[i][0] = "missing_actual"
    for j in range(1, actual_count + 1):
        dp[0][j] = dp[0][j - 1] + gap_penalty
        step[0][j] = "extra_actual"

    for i in range(1, expected_count + 1):
        for j in range(1, actual_count + 1):
            match_gain = scores[i - 1][j - 1] - min_match_score
            candidates = (
                (dp[i - 1][j - 1] + match_gain, "matched"),
                (dp[i - 1][j] + gap_penalty, "missing_actual"),
                (dp[i][j - 1] + gap_penalty, "extra_actual"),
            )
            dp[i][j], step[i][j] = max(candidates, key=lambda item: item[0])

    matches: List[SlideMatch] = []
    i = expected_count
    j = actual_count
    while i > 0 or j > 0:
        current = step[i][j]
        if current == "matched":
            score = scores[i - 1][j - 1]
            # Trust the DP: if it chose to match, we compare these two slides.
            # A below-threshold score is flagged 'low_confidence_match' so the
            # user can see it, but it is never split into missing+extra — that
            # would throw away the alignment the DP just computed.
            status = "matched" if score >= min_match_score else "low_confidence_match"
            matches.append(SlideMatch(i - 1, j - 1, score, status))
            i -= 1
            j -= 1
        elif current == "missing_actual" or j == 0:
            matches.append(SlideMatch(i - 1, None, None, "missing_actual"))
            i -= 1
        else:
            matches.append(SlideMatch(None, j - 1, None, "extra_actual"))
            j -= 1

    matches.reverse()
    return matches


def align_slide_pages(
    expected_pages: Sequence[Image.Image],
    actual_pages: Sequence[Image.Image],
    min_match_score: float,
    expected_texts: Optional[Sequence[str]] = None,
    actual_texts: Optional[Sequence[str]] = None,
    gap_penalty: float = DEFAULT_GAP_PENALTY,
) -> List[SlideMatch]:
    """Align two slide sequences by visual (and, when available, text) similarity."""
    expected_sigs = [
        compute_signature(page, _text_at(expected_texts, i))
        for i, page in enumerate(expected_pages)
    ]
    actual_sigs = [
        compute_signature(page, _text_at(actual_texts, j))
        for j, page in enumerate(actual_pages)
    ]
    scores = build_score_matrix(expected_sigs, actual_sigs)
    return align_from_score_matrix(
        scores, len(expected_pages), len(actual_pages), min_match_score, gap_penalty
    )


def index_slide_pages(expected_pages: Sequence[Image.Image], actual_pages: Sequence[Image.Image]) -> List[SlideMatch]:
    page_count = max(len(expected_pages), len(actual_pages))
    matches: List[SlideMatch] = []
    for index in range(page_count):
        expected_index = index if index < len(expected_pages) else None
        actual_index = index if index < len(actual_pages) else None
        if expected_index is not None and actual_index is not None:
            matches.append(SlideMatch(expected_index, actual_index, None, "same_index"))
        elif expected_index is not None:
            matches.append(SlideMatch(expected_index, None, None, "missing_actual"))
        else:
            matches.append(SlideMatch(None, actual_index, None, "extra_actual"))
    return matches


def _text_at(texts: Optional[Sequence[str]], index: int) -> Optional[str]:
    if texts is None or index >= len(texts):
        return None
    return texts[index]


def _maybe_load_slide_texts(path: Path) -> Optional[List[str]]:
    """Per-slide text for PPTX inputs; None for images/PPT where text is unavailable."""
    if path.suffix.lower() != PPTX_EXT:
        return None
    try:
        return extract_slide_texts(path)
    except Exception:
        return None


def _maybe_load_slide_titles(path: Path) -> Optional[List[str]]:
    """Per-slide title text for PPTX inputs; None when not available."""
    if path.suffix.lower() != PPTX_EXT:
        return None
    try:
        return extract_slide_titles(path)
    except Exception:
        return None


def write_similarity_matrix(
    output_dir: Path,
    scores: Sequence[Sequence[float]],
    matches: Sequence[SlideMatch],
    min_match_score: float,
    text_used: bool,
) -> None:
    """Dump the full N×M similarity matrix and the chosen alignment as JSON.

    This is the matcher's audit trail: the user can see *why* two slides were
    paired — or why a slide was left unmatched — instead of trusting the result
    blindly.
    """
    expected_count = len(scores)
    actual_count = len(scores[0]) if scores and scores[0] else 0

    best_match_per_expected = []
    for i in range(expected_count):
        row = scores[i]
        if not row:
            continue
        best_j = max(range(len(row)), key=lambda j: row[j])
        best_match_per_expected.append(
            {
                "expected_slide": i + 1,
                "best_actual_slide": best_j + 1,
                "score": round(row[best_j], 6),
            }
        )

    payload = {
        "expected_count": expected_count,
        "actual_count": actual_count,
        "min_match_score": min_match_score,
        "text_signal_used": text_used,
        "scores": [[round(value, 6) for value in row] for row in scores],
        "best_match_per_expected": best_match_per_expected,
        "alignment": [
            {
                "expected_slide": match.expected_index + 1 if match.expected_index is not None else None,
                "actual_slide": match.actual_index + 1 if match.actual_index is not None else None,
                "score": round(match.score, 6) if match.score is not None else None,
                "status": match.status,
            }
            for match in matches
        ],
    }
    (output_dir / "similarity_matrix.json").write_text(
        json.dumps(payload, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )


def compare_pixels(
    args: argparse.Namespace, output_dir: Path
) -> Tuple[List[PageDiff], List[SlideMatch]]:
    """Render & diff slides; returns per-pair diffs *and* the alignment used.

    Returning the alignment lets the object-comparison pass reuse the same slide
    pairing, so a reordered slide is not mistakenly diffed against the wrong
    one.
    """
    expected_path = Path(args.expected)
    actual_path = Path(args.actual)

    temp_root = Path(tempfile.mkdtemp(prefix="report_diff_render_"))
    try:
        expected_pages = load_pages(expected_path, args.dpi, temp_root / "expected")
        actual_pages = load_pages(actual_path, args.dpi, temp_root / "actual")
    finally:
        shutil.rmtree(str(temp_root), ignore_errors=True)

    min_match_score = getattr(args, "min_match_score", 0.82)
    gap_penalty = getattr(args, "gap_penalty", DEFAULT_GAP_PENALTY)
    align_slides = bool(getattr(args, "align_slides", False))

    if align_slides:
        expected_texts = _maybe_load_slide_texts(expected_path)
        actual_texts = _maybe_load_slide_texts(actual_path)
        expected_titles = _maybe_load_slide_titles(expected_path)
        actual_titles = _maybe_load_slide_titles(actual_path)
        # Text/title are only usable matching signals when *both* reports expose them.
        if expected_texts is None or actual_texts is None:
            expected_texts = actual_texts = None
        if expected_titles is None or actual_titles is None:
            expected_titles = actual_titles = None

        expected_sigs = [
            compute_signature(page, _text_at(expected_texts, i), _text_at(expected_titles, i))
            for i, page in enumerate(expected_pages)
        ]
        actual_sigs = [
            compute_signature(page, _text_at(actual_texts, j), _text_at(actual_titles, j))
            for j, page in enumerate(actual_pages)
        ]
        scores = build_score_matrix(expected_sigs, actual_sigs)
        matches = align_from_score_matrix(
            scores, len(expected_pages), len(actual_pages), min_match_score, gap_penalty
        )
        write_similarity_matrix(
            output_dir, scores, matches, min_match_score, text_used=expected_texts is not None
        )
    else:
        matches = index_slide_pages(expected_pages, actual_pages)

    results = []
    blank = Image.new("RGB", (1, 1), (255, 255, 255))

    for compare_index, match in enumerate(matches, start=1):
        expected = expected_pages[match.expected_index] if match.expected_index is not None else blank
        actual = actual_pages[match.actual_index] if match.actual_index is not None else blank
        results.append(
            compare_page(
                expected=expected,
                actual=actual,
                page_number=compare_index,
                expected_page=match.expected_index + 1 if match.expected_index is not None else None,
                actual_page=match.actual_index + 1 if match.actual_index is not None else None,
                match_score=match.score,
                match_status=match.status,
                output_dir=output_dir,
                threshold=args.threshold,
                allowed_percent=args.allowed_percent,
                highlight_color=args.highlight_color,
                alpha=args.alpha,
            )
        )

    return results, list(matches)


def hash_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def emu_to_points(value: int) -> float:
    return round(float(value) / EMU_PER_POINT, 3)


def shape_to_record(shape: Any) -> Dict[str, Any]:
    record = {
        "name": safe_text(getattr(shape, "name", "")),
        "shape_type": safe_text(getattr(shape, "shape_type", "")),
        "left_pt": emu_to_points(getattr(shape, "left", 0)),
        "top_pt": emu_to_points(getattr(shape, "top", 0)),
        "width_pt": emu_to_points(getattr(shape, "width", 0)),
        "height_pt": emu_to_points(getattr(shape, "height", 0)),
        "text": "",
        "image_sha256": "",
    }

    if getattr(shape, "has_text_frame", False):
        record["text"] = safe_text(shape.text_frame.text)

    if hasattr(shape, "image"):
        try:
            record["image_sha256"] = hash_bytes(shape.image.blob)
        except Exception:
            record["image_sha256"] = ""

    return record


def extract_pptx_objects(path: Path) -> List[List[Dict[str, Any]]]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "PPTX object comparison requires python-pptx. Install it with: pip install python-pptx"
        ) from exc

    presentation = Presentation(str(path))
    slides = []
    for slide in presentation.slides:
        slides.append([shape_to_record(shape) for shape in slide.shapes])
    return slides


def extract_slide_texts(path: Path) -> List[str]:
    """Concatenated text content per slide — a high-confidence matching signal.

    Two slides that share a title (e.g. an analysis-conditions page) are almost
    certainly the same slide across report versions, even if their charts moved.
    """
    slides = extract_pptx_objects(path)
    return [
        "\n".join(record["text"] for record in shapes if record.get("text"))
        for shapes in slides
    ]


def extract_slide_titles(path: Path) -> List[str]:
    """Per-slide title text — the strongest template signal for technical reports.

    Uses the PPTX title placeholder when present; otherwise falls back to the
    topmost meaningfully-long text shape on the slide.  Two slides labelled
    "解析条件 - 燃焼圧" (with different test cases below) match almost perfectly
    on title alone, which is exactly the kind of cross-version pairing the body
    text — with its differing values — would muddy.
    """
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError(
            "Title extraction requires python-pptx. Install it with: pip install python-pptx"
        ) from exc

    titles: List[str] = []
    presentation = Presentation(str(path))
    for slide in presentation.slides:
        title = ""
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title = (title_shape.text_frame.text or "").strip()

        if not title:
            # Fallback: topmost text shape with at least 4 characters — skips
            # short page numbers and "Confidential" stamps at the corner.
            candidates = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = (shape.text_frame.text or "").strip()
                if len(text) < 4:
                    continue
                top = getattr(shape, "top", 0) or 0
                candidates.append((top, text))
            if candidates:
                title = min(candidates, key=lambda candidate: candidate[0])[1]

        titles.append(title)
    return titles


def align_pptx_by_text(
    expected_path: Path,
    actual_path: Path,
    min_match_score: float,
    gap_penalty: float = DEFAULT_GAP_PENALTY,
    output_dir: Optional[Path] = None,
) -> List[SlideMatch]:
    """Align two PPTX files using title + body text similarity (no rendering).

    For object-only comparisons we have no pixel signatures to work with, but
    PPTX gives us per-slide titles and full body text — usually a stronger
    matching signal than the visual path anyway.  Title bigram Jaccard
    dominates; the body provides a soft tie-breaker.  If ``output_dir`` is
    given, the score matrix is exported for audit just like the visual path.
    """
    expected_texts = extract_slide_texts(expected_path)
    actual_texts = extract_slide_texts(actual_path)
    try:
        expected_titles = extract_slide_titles(expected_path)
        actual_titles = extract_slide_titles(actual_path)
    except Exception:
        expected_titles = ["" for _ in expected_texts]
        actual_titles = ["" for _ in actual_texts]

    def pair_score(i: int, j: int) -> float:
        title_sim = text_similarity(
            _text_at(expected_titles, i) or "",
            _text_at(actual_titles, j) or "",
        )
        body_sim = text_similarity(expected_texts[i], actual_texts[j])
        if title_sim >= 0.85:
            return 0.85 * (0.70 * title_sim + 0.30 * body_sim) + 0.15 * body_sim
        if title_sim >= 0.55:
            return 0.65 * (0.70 * title_sim + 0.30 * body_sim) + 0.35 * body_sim
        return 0.45 * title_sim + 0.55 * body_sim

    expected_count = len(expected_texts)
    actual_count = len(actual_texts)
    scores = [
        [pair_score(i, j) for j in range(actual_count)]
        for i in range(expected_count)
    ]
    matches = align_from_score_matrix(
        scores, expected_count, actual_count, min_match_score, gap_penalty
    )
    if output_dir is not None:
        write_similarity_matrix(output_dir, scores, matches, min_match_score, text_used=True)
    return matches


def compare_objects(
    args: argparse.Namespace,
    output_dir: Path,
    matches: Optional[Sequence[SlideMatch]] = None,
) -> List[ObjectDiff]:
    """Diff shape-level data slide by slide, using a shared alignment when given.

    If ``matches`` is provided (typically by an upstream pixel-comparison pass)
    those exact slide pairings are reused, so reordered slides are diffed
    against their real counterpart rather than against whatever happens to sit
    at the same index. When ``matches`` is None this function computes its own
    alignment: text-based for PPTX when ``--align-slides`` is on, otherwise the
    historical index-based pairing.
    """
    expected_path = Path(args.expected)
    actual_path = Path(args.actual)
    if expected_path.suffix.lower() != PPTX_EXT or actual_path.suffix.lower() != PPTX_EXT:
        raise ValueError("Object comparison currently supports .pptx files only.")

    expected_slides = extract_pptx_objects(expected_path)
    actual_slides = extract_pptx_objects(actual_path)

    if matches is None:
        if bool(getattr(args, "align_slides", False)):
            matches = align_pptx_by_text(
                expected_path,
                actual_path,
                getattr(args, "min_match_score", 0.82),
                getattr(args, "gap_penalty", DEFAULT_GAP_PENALTY),
                output_dir=output_dir,
            )
        else:
            matches = index_slide_pages(expected_slides, actual_slides)

    diffs: List[ObjectDiff] = []
    for pair_index, match in enumerate(matches, start=1):
        expected_shapes = (
            expected_slides[match.expected_index]
            if match.expected_index is not None and match.expected_index < len(expected_slides)
            else []
        )
        actual_shapes = (
            actual_slides[match.actual_index]
            if match.actual_index is not None and match.actual_index < len(actual_slides)
            else []
        )
        expected_slide_no = match.expected_index + 1 if match.expected_index is not None else None
        actual_slide_no = match.actual_index + 1 if match.actual_index is not None else None
        object_count = max(len(expected_shapes), len(actual_shapes))

        for object_index in range(object_count):
            expected = expected_shapes[object_index] if object_index < len(expected_shapes) else {}
            actual = actual_shapes[object_index] if object_index < len(actual_shapes) else {}
            fields = sorted(set(expected.keys()) | set(actual.keys()))
            for field in fields:
                if expected.get(field) != actual.get(field):
                    diffs.append(
                        ObjectDiff(
                            pair_index=pair_index,
                            expected_slide=expected_slide_no,
                            actual_slide=actual_slide_no,
                            object_index=object_index + 1,
                            field=field,
                            expected=expected.get(field),
                            actual=actual.get(field),
                        )
                    )

    object_summary = output_dir / "object_summary.json"
    object_summary.write_text(
        json.dumps([asdict(diff) for diff in diffs], indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    return diffs


def get_package_root() -> Path:
    return Path(__file__).resolve().parent


def resolve_output_dir(output_dir: str) -> Path:
    package_root = get_package_root()
    requested = Path(output_dir)
    if not requested.is_absolute():
        requested = package_root / requested

    resolved = requested.resolve()
    try:
        resolved.relative_to(package_root)
    except ValueError:
        raise ValueError(
            "Output directory must stay inside the tool package folder: {}".format(
                package_root
            )
        )

    return resolved


def validate_inputs(args: argparse.Namespace) -> None:
    expected_path = Path(args.expected)
    actual_path = Path(args.actual)

    if not expected_path.exists():
        raise FileNotFoundError("Expected report not found: {}".format(expected_path))
    if not actual_path.exists():
        raise FileNotFoundError("Actual report not found: {}".format(actual_path))

    if args.threshold < 0 or args.threshold > 255:
        raise ValueError("--threshold must be between 0 and 255")
    if getattr(args, "min_match_score", 0.82) < 0 or getattr(args, "min_match_score", 0.82) > 1:
        raise ValueError("--min-match-score must be between 0 and 1")
    if args.allowed_percent < 0 or args.allowed_percent > 100:
        raise ValueError("--allowed-percent must be between 0 and 100")
    if args.alpha < 0 or args.alpha > 255:
        raise ValueError("--alpha must be between 0 and 255")
    if args.dpi <= 0:
        raise ValueError("--dpi must be greater than 0")


def write_summary(output_dir: Path, pixel_results: Sequence[PageDiff], object_diffs: Sequence[ObjectDiff]) -> None:
    payload = {
        "pixel_results": [asdict(result) for result in pixel_results],
        "object_diffs": [asdict(diff) for diff in object_diffs],
    }
    (output_dir / "summary.json").write_text(
        json.dumps(payload, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )


def print_summary(pixel_results: Iterable[PageDiff], object_diffs: Sequence[ObjectDiff]) -> None:
    total_pages = 0
    failed_pages = 0

    for result in pixel_results:
        total_pages += 1
        if not result.passed:
            failed_pages += 1

        status = "PASS" if result.passed else "FAIL"
        print(
            "[{}] pair={} expected={} actual={} match={} score={} diff={:.6f}% pixels={}/{} bbox={} overlay={}".format(
                status,
                result.page,
                result.expected_page,
                result.actual_page,
                result.match_status,
                result.match_score,
                result.difference_percent,
                result.different_pixels,
                result.compared_pixels,
                result.bbox,
                result.output_overlay,
            )
        )

    if total_pages:
        print("Pages compared: {}; failed pages: {}".format(total_pages, failed_pages))

    if object_diffs:
        print("Object differences: {}".format(len(object_diffs)))
        for diff in object_diffs[:20]:
            print(
                "[OBJECT] pair={} expected_slide={} actual_slide={} object={} field={}".format(
                    diff.pair_index,
                    diff.expected_slide if diff.expected_slide is not None else "-",
                    diff.actual_slide if diff.actual_slide is not None else "-",
                    diff.object_index,
                    diff.field,
                )
            )
        if len(object_diffs) > 20:
            print("... {} more object differences in object_summary.json".format(len(object_diffs) - 20))
    elif total_pages == 0:
        print("Object differences: 0")


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Compare report files by pixels and/or PowerPoint objects."
    )
    parser.add_argument("expected", help="Baseline report image/PPT/PPTX")
    parser.add_argument("actual", help="Report image/PPT/PPTX to check")
    parser.add_argument(
        "-o",
        "--output-dir",
        default=DEFAULT_OUTPUT_DIR,
        help="Directory for overlay images, masks, and JSON summaries",
    )
    parser.add_argument(
        "--mode",
        choices=("pixel", "object", "both"),
        default="pixel",
        help="pixel renders pages/slides; object compares PPTX shapes; both runs both checks",
    )
    parser.add_argument(
        "--align-slides",
        action="store_true",
        help="Auto-match slides by visual similarity before comparing",
    )
    parser.add_argument(
        "--min-match-score",
        type=float,
        default=0.82,
        help="Minimum visual similarity score for auto slide matching, 0.0 to 1.0",
    )
    parser.add_argument(
        "--gap-penalty",
        type=float,
        default=DEFAULT_GAP_PENALTY,
        help="Cost of leaving a slide unmatched during auto alignment; a more "
        "negative value makes the matcher less willing to pair dissimilar slides",
    )
    parser.add_argument(
        "-t",
        "--threshold",
        type=int,
        default=0,
        help="Ignore per-channel pixel differences at or below this value",
    )
    parser.add_argument(
        "--allowed-percent",
        type=float,
        default=0.0,
        help="Maximum allowed percent of different pixels before a page fails",
    )
    parser.add_argument(
        "--highlight-color",
        type=parse_color,
        default=(255, 0, 0),
        help="Overlay color as R,G,B",
    )
    parser.add_argument(
        "--alpha",
        type=int,
        default=51,
        help="Highlight fill opacity from 0 to 255; 51 is 20 percent",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=150,
        help="DPI used when rendering PowerPoint slides",
    )
    return parser


def main(argv: Optional[List[str]] = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)

    try:
        validate_inputs(args)
        output_dir = resolve_output_dir(args.output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

        pixel_results: List[PageDiff] = []
        object_diffs: List[ObjectDiff] = []
        matches: Optional[List[SlideMatch]] = None

        if args.mode in ("pixel", "both"):
            pixel_results, matches = compare_pixels(args, output_dir)
        if args.mode in ("object", "both"):
            object_diffs = compare_objects(args, output_dir, matches=matches)

        write_summary(output_dir, pixel_results, object_diffs)
    except Exception as exc:
        print("Error: {}".format(exc), file=sys.stderr)
        return 2

    print_summary(pixel_results, object_diffs)
    pixels_passed = all(result.passed for result in pixel_results)
    objects_passed = len(object_diffs) == 0
    return 0 if pixels_passed and objects_passed else 1


if __name__ == "__main__":
    raise SystemExit(main())






